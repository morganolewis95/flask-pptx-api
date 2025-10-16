from flask import Flask, request, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import tempfile
import os
import requests
from io import BytesIO

app = Flask(__name__)
CORS(app)

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    data = request.json
    slides = data.get("slides", [])

    prs = Presentation()

    for slide_data in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Background color
        bg_color = slide_data.get("background_color")
        if bg_color:
            fill = slide.background.fill
            fill.solid()
            rgb = tuple(int(bg_color.lstrip('#')[i:i+2], 16) for i in (0, 2 ,4))
            fill.fore_color.rgb = RGBColor(*rgb)

        # Title formatting
        title_shape = slide.shapes.title
        title_text = slide_data.get("title", "Untitled")
        title_shape.text = title_text
        if title_shape.text_frame.paragraphs:
            title_run = title_shape.text_frame.paragraphs[0].runs[0]
            title_run.font.bold = True
            title_run.font.size = Pt(32)

        # Bullet points
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()
        bullets = slide_data.get("bullets", [])
        for bullet in bullets:
            p = content_frame.add_paragraph()
            p.text = bullet
            p.level = 0

        # Image
        image_url = slide_data.get("image_url")
        if image_url:
            try:
                img_resp = requests.get(image_url)
                if img_resp.status_code == 200:
                    image_stream = BytesIO(img_resp.content)
                    slide.shapes.add_picture(image_stream, Inches(5), Inches(2), width=Inches(4.5))
            except Exception as e:
                print(f"Image download failed: {e}")

    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp_file.name)

    return send_file(temp_file.name, as_attachment=True, download_name="generated_presentation.pptx")

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port, debug=True)
